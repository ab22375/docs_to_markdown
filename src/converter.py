#!/usr/bin/env python3
"""
Document to Markdown/JSON converter
Recursively scans directories for PDF, PPTX, DOCX files and converts them to .md and .json
"""

# Force CPU usage to avoid MPS issues on macOS - must be before any imports
import os
os.environ["PYTORCH_MPS_DISABLED"] = "1"
os.environ["PYTORCH_DEVICE"] = "cpu"
os.environ["CUDA_VISIBLE_DEVICES"] = ""

import concurrent.futures
import json
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import click
import torch

# Disable MPS backend if available
if hasattr(torch.backends, 'mps') and torch.backends.mps.is_available():
    torch.backends.mps.enabled = False
from docx import Document
from marker.converters.pdf import PdfConverter
from marker.models import create_model_dict
from marker.output import text_from_rendered
from pptx import Presentation
from rich.console import Console
from rich.panel import Panel
from rich.progress import BarColumn, Progress, SpinnerColumn, TextColumn
from rich.table import Table

from .ocr import OCRManager

console = Console()


@dataclass
class ConversionResult:
    """Store conversion results"""

    source_path: str
    markdown_path: str
    json_path: str
    status: str
    error: Optional[str] = None
    metadata: Optional[Dict] = None


class DocumentConverter:
    """Main converter class for all document types"""

    def __init__(
        self, output_dir: Optional[Path] = None, parallel: bool = True, base_input_path: Optional[Path] = None, enable_ocr: bool = False
    ):
        self.output_dir = output_dir
        self.parallel = parallel
        # Force CPU to avoid MPS issues on macOS with marker-pdf
        self.device = torch.device("cpu")
        self.pdf_converter = None
        self.base_input_path = base_input_path
        self.ocr_manager = OCRManager(self.device, enable_ocr=enable_ocr)

    def _load_pdf_converter(self):
        """Lazy load PDF converter"""
        if self.pdf_converter is None:
            self.pdf_converter = PdfConverter(
                artifact_dict=create_model_dict(),
            )

    def _get_output_paths(self, input_path: Path) -> Tuple[Path, Path]:
        """Generate output paths for markdown and json files in separate subdirectories"""
        if self.output_dir:
            # If we have a base input path (directory was passed), preserve structure
            if self.base_input_path and self.base_input_path.is_dir():
                try:
                    # Get the relative path from the base input directory
                    relative_path = input_path.relative_to(self.base_input_path)
                    # Create separate md and json subdirectories
                    md_base = self.output_dir / "md" / relative_path.parent / input_path.stem
                    json_base = self.output_dir / "json" / relative_path.parent / input_path.stem
                except ValueError:
                    # If the file is not under base_input_path, just use the filename
                    md_base = self.output_dir / "md" / input_path.stem
                    json_base = self.output_dir / "json" / input_path.stem
            else:
                # Single file was passed, just use the filename
                md_base = self.output_dir / "md" / input_path.stem
                json_base = self.output_dir / "json" / input_path.stem
        else:
            # When no output dir specified, put files next to original
            base = input_path.parent / input_path.stem
            return base.with_suffix(".md"), base.with_suffix(".json")

        return md_base.with_suffix(".md"), json_base.with_suffix(".json")

    def convert_pdf(self, pdf_path: Path) -> ConversionResult:
        """Convert PDF to markdown using marker-pdf with OCR fallback for scanned PDFs"""
        try:
            # Check if OCR is needed for this PDF
            ocr_result = self.ocr_manager.process_if_needed(pdf_path)

            if ocr_result:
                # OCR was used
                full_text, ocr_metadata = ocr_result
                metadata = {"ocr_used": True, **ocr_metadata}
                images_count = ocr_metadata.get("page_count", 0)
            else:
                # Use standard marker-pdf conversion
                self._load_pdf_converter()
                rendered = self.pdf_converter(str(pdf_path))
                full_text, _, images = text_from_rendered(rendered)
                metadata = {"ocr_used": False}
                images_count = len(images) if images else 0

            # Get output paths
            md_path, json_path = self._get_output_paths(pdf_path)
            md_path.parent.mkdir(parents=True, exist_ok=True)
            json_path.parent.mkdir(parents=True, exist_ok=True)

            # Save markdown
            md_path.write_text(full_text, encoding="utf-8")

            # Save JSON with metadata
            json_data = {
                "source": str(pdf_path),
                "type": "pdf",
                "content": full_text,
                "metadata": metadata,
                "images": images_count,
            }
            json_path.write_text(json.dumps(json_data, indent=2), encoding="utf-8")

            return ConversionResult(
                source_path=str(pdf_path),
                markdown_path=str(md_path),
                json_path=str(json_path),
                status="success",
                metadata=metadata,
            )

        except Exception as e:
            return ConversionResult(
                source_path=str(pdf_path), markdown_path="", json_path="", status="error", error=str(e)
            )

    def convert_docx(self, docx_path: Path) -> ConversionResult:
        """Convert DOCX to markdown"""
        try:
            doc = Document(docx_path)

            # Extract text with basic formatting
            markdown_lines = []

            for para in doc.paragraphs:
                if para.style.name.startswith("Heading"):
                    level = int(para.style.name[-1]) if para.style.name[-1].isdigit() else 1
                    markdown_lines.append(f"{'#' * level} {para.text}")
                elif para.text.strip():
                    # Handle basic formatting
                    text = para.text
                    for run in para.runs:
                        if run.bold:
                            text = text.replace(run.text, f"**{run.text}**", 1)
                        elif run.italic:
                            text = text.replace(run.text, f"*{run.text}*", 1)
                    markdown_lines.append(text)
                else:
                    markdown_lines.append("")

            full_text = "\n\n".join(markdown_lines)

            # Get output paths
            md_path, json_path = self._get_output_paths(docx_path)
            md_path.parent.mkdir(parents=True, exist_ok=True)
            json_path.parent.mkdir(parents=True, exist_ok=True)

            # Save markdown
            md_path.write_text(full_text, encoding="utf-8")

            # Save JSON
            json_data = {
                "source": str(docx_path),
                "type": "docx",
                "content": full_text,
                "metadata": {
                    "author": doc.core_properties.author or "",
                    "created": doc.core_properties.created.isoformat() if doc.core_properties.created else "",
                    "modified": doc.core_properties.modified.isoformat() if doc.core_properties.modified else "",
                    "title": doc.core_properties.title or "",
                },
            }
            json_path.write_text(json.dumps(json_data, indent=2), encoding="utf-8")

            return ConversionResult(
                source_path=str(docx_path),
                markdown_path=str(md_path),
                json_path=str(json_path),
                status="success",
                metadata=json_data["metadata"],
            )

        except Exception as e:
            return ConversionResult(
                source_path=str(docx_path), markdown_path="", json_path="", status="error", error=str(e)
            )

    def convert_pptx(self, pptx_path: Path) -> ConversionResult:
        """Convert PPTX to markdown"""
        try:
            prs = Presentation(pptx_path)

            markdown_lines = []

            for i, slide in enumerate(prs.slides, 1):
                markdown_lines.append(f"# Slide {i}")

                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Check if it's a title
                        if shape == slide.shapes.title:
                            markdown_lines.append(f"## {shape.text.strip()}")
                        else:
                            markdown_lines.append(shape.text.strip())

                # Add notes if present
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                    markdown_lines.append("\n**Notes:**")
                    markdown_lines.append(slide.notes_slide.notes_text_frame.text.strip())

                markdown_lines.append("")  # Add spacing between slides

            full_text = "\n\n".join(markdown_lines)

            # Get output paths
            md_path, json_path = self._get_output_paths(pptx_path)
            md_path.parent.mkdir(parents=True, exist_ok=True)
            json_path.parent.mkdir(parents=True, exist_ok=True)

            # Save markdown
            md_path.write_text(full_text, encoding="utf-8")

            # Save JSON
            json_data = {
                "source": str(pptx_path),
                "type": "pptx",
                "content": full_text,
                "metadata": {
                    "slide_count": len(prs.slides),
                    "title": prs.core_properties.title or "",
                    "author": prs.core_properties.author or "",
                    "created": prs.core_properties.created.isoformat() if prs.core_properties.created else "",
                    "modified": prs.core_properties.modified.isoformat() if prs.core_properties.modified else "",
                },
            }
            json_path.write_text(json.dumps(json_data, indent=2), encoding="utf-8")

            return ConversionResult(
                source_path=str(pptx_path),
                markdown_path=str(md_path),
                json_path=str(json_path),
                status="success",
                metadata=json_data["metadata"],
            )

        except Exception as e:
            return ConversionResult(
                source_path=str(pptx_path), markdown_path="", json_path="", status="error", error=str(e)
            )

    def find_documents(self, path: Path, skip_pdf: bool = False) -> List[Path]:
        """Recursively find all supported documents"""
        supported_extensions = {".pdf", ".docx", ".pptx"}
        
        if skip_pdf:
            supported_extensions.discard(".pdf")
            console.print("[yellow]Note: PDF files are skipped due to compatibility issues[/yellow]")
        
        documents = []

        if path.is_file():
            if path.suffix.lower() in supported_extensions:
                documents.append(path)
        else:
            for ext in supported_extensions:
                documents.extend(path.rglob(f"*{ext}"))

        return sorted(documents)

    def convert_document(self, doc_path: Path) -> ConversionResult:
        """Convert a single document based on its type"""
        suffix = doc_path.suffix.lower()

        if suffix == ".pdf":
            return self.convert_pdf(doc_path)
        elif suffix == ".docx":
            return self.convert_docx(doc_path)
        elif suffix == ".pptx":
            return self.convert_pptx(doc_path)
        else:
            return ConversionResult(
                source_path=str(doc_path),
                markdown_path="",
                json_path="",
                status="error",
                error=f"Unsupported file type: {suffix}",
            )

    def convert_all(self, path: Path, skip_pdf: bool = False) -> List[ConversionResult]:
        """Convert all documents found in path"""
        documents = self.find_documents(path, skip_pdf=skip_pdf)

        if not documents:
            console.print("[yellow]No supported documents found!")
            return []

        console.print(f"[green]Found {len(documents)} documents to convert")

        results = []

        with Progress(
            SpinnerColumn(),
            TextColumn("[progress.description]{task.description}"),
            BarColumn(),
            TextColumn("[progress.percentage]{task.percentage:>3.0f}%"),
            console=console,
        ) as progress:
            task = progress.add_task("Converting documents...", total=len(documents))

            if self.parallel and len(documents) > 1:
                with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
                    future_to_doc = {executor.submit(self.convert_document, doc): doc for doc in documents}

                    for future in concurrent.futures.as_completed(future_to_doc):
                        result = future.result()
                        results.append(result)
                        progress.update(task, advance=1)

                        if result.status == "success":
                            console.print(f"[green]✓[/green] {result.source_path}")
                        else:
                            console.print(f"[red]✗[/red] {result.source_path}: {result.error}")
            else:
                for doc in documents:
                    result = self.convert_document(doc)
                    results.append(result)
                    progress.update(task, advance=1)

                    if result.status == "success":
                        console.print(f"[green]✓[/green] {result.source_path}")
                    else:
                        console.print(f"[red]✗[/red] {result.source_path}: {result.error}")

        return results


@click.command()
@click.argument("path", type=click.Path(exists=True, path_type=Path))
@click.option("--output-dir", "-o", type=click.Path(path_type=Path), help="Output directory (default: same as source)")
@click.option("--no-parallel", is_flag=True, help="Disable parallel processing")
@click.option("--summary", "-s", is_flag=True, help="Show conversion summary table")
@click.option("--enable-ocr", is_flag=True, help="Enable OCR for scanned PDFs (experimental)")
@click.option("--skip-pdf", is_flag=True, help="Skip PDF files (recommended due to compatibility issues)")
def main(path: Path, output_dir: Optional[Path], no_parallel: bool, summary: bool, enable_ocr: bool, skip_pdf: bool):
    """
    Convert PDF, DOCX, and PPTX files to Markdown and JSON.

    PATH can be a single file or directory to scan recursively.
    """
    console.print(
        Panel.fit(f"[bold blue]Document to Markdown Converter[/bold blue]\nProcessing: {path}", border_style="blue")
    )

    
    converter = DocumentConverter(
        output_dir=output_dir, parallel=not no_parallel, base_input_path=path if path.is_dir() else None, enable_ocr=enable_ocr
    )

    results = converter.convert_all(path, skip_pdf=skip_pdf)

    if summary and results:
        # Create summary table
        table = Table(title="Conversion Summary")
        table.add_column("File", style="cyan")
        table.add_column("Status", style="green")
        table.add_column("Method", style="magenta")
        table.add_column("Output", style="yellow")

        success_count = 0
        ocr_count = 0
        for result in results:
            status = "✓ Success" if result.status == "success" else f"✗ Error: {result.error}"
            status_style = "green" if result.status == "success" else "red"

            # Determine conversion method
            method = "Standard"
            if result.status == "success" and result.metadata:
                if result.metadata.get("ocr_used", False):
                    method = f"OCR ({result.metadata.get('ocr_engine', 'unknown')})"
                    ocr_count += 1

            if result.status == "success":
                success_count += 1
                output = f"{result.markdown_path}"
            else:
                output = "N/A"

            table.add_row(
                os.path.basename(result.source_path), f"[{status_style}]{status}[/{status_style}]", method, output
            )

        console.print(table)

        # Enhanced summary with OCR info
        failed_count = len(results) - success_count
        summary_parts = [f"Total: {len(results)}", f"Success: {success_count}", f"Failed: {failed_count}"]
        if ocr_count > 0:
            summary_parts.append(f"OCR Used: {ocr_count}")

        console.print(f"\n[bold]{' | '.join(summary_parts)}[/bold]")


if __name__ == "__main__":
    main()
